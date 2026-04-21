"""Generate a professional PowerPoint presentation for the ANN Character Recognition project."""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Colour palette ──────────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
ACCENT_BLUE   = RGBColor(0x38, 0x7A, 0xF5)   # Bright blue
ACCENT_TEAL   = RGBColor(0x2A, 0x9D, 0x8F)   # Teal
ACCENT_ORANGE = RGBColor(0xF9, 0x73, 0x16)   # Orange
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xA0, 0xAE, 0xC0)
DARK_GRAY     = RGBColor(0x64, 0x74, 0x8B)
CARD_BG       = RGBColor(0x1E, 0x29, 0x3B)   # Slightly lighter navy for cards
SECTION_BG    = RGBColor(0x15, 0x20, 0x38)   # Section dividers


SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

OUTPUT_DIR = Path("outputs")
PPTX_PATH  = Path("ANN_Character_Recognition_Presentation.pptx")


# ── Helpers ─────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT_TEAL, font_name="Segoe UI",
                    line_spacing=Pt(28)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = line_spacing
        p.level = 0

        # Bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '●'})
        # Remove existing bullets
        for existing in pPr.findall(qn('a:buChar')):
            pPr.remove(existing)
        for existing in pPr.findall(qn('a:buNone')):
            pPr.remove(existing)
        pPr.append(buChar)

        # Bullet color
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': f'{bullet_color}'})
        buClr.append(srgbClr)
        for existing in pPr.findall(qn('a:buClr')):
            pPr.remove(existing)
        pPr.append(buClr)

    return txBox


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add image if it exists, otherwise add a placeholder box."""
    if Path(img_path).exists():
        kwargs = {"image_file": str(img_path), "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(**kwargs)
    else:
        shape = add_shape_rect(slide, left, top,
                               width or Inches(4), height or Inches(3),
                               CARD_BG, DARK_GRAY, Pt(1))
        shape.text_frame.paragraphs[0].text = f"[Image: {Path(img_path).name}]"
        shape.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
        shape.text_frame.paragraphs[0].font.size = Pt(12)


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_stat_card(slide, left, top, value, label, accent=ACCENT_BLUE):
    card_w = Inches(2.6)
    card_h = Inches(1.6)
    add_shape_rect(slide, left, top, card_w, card_h, CARD_BG)
    add_accent_line(slide, left + Inches(0.3), top + Inches(0.15), Inches(2.0), accent)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.3), Inches(2.0), Inches(0.7),
                 value, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.95), Inches(2.0), Inches(0.5),
                 label, font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)


# ── Slide builders ──────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Decorative accent line at top
    add_accent_line(slide, Inches(1.5), Inches(1.5), Inches(2.5), ACCENT_TEAL)

    add_text_box(slide, Inches(1.5), Inches(1.7), Inches(10), Inches(1.2),
                 "Handwritten Character Recognition", font_size=40, color=WHITE, bold=True)

    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
                 "Using a From-Scratch NumPy Multi-Layer Perceptron",
                 font_size=22, color=ACCENT_BLUE, bold=False)

    add_accent_line(slide, Inches(1.5), Inches(3.7), Inches(5), CARD_BG)

    add_text_box(slide, Inches(1.5), Inches(4.1), Inches(8), Inches(0.5),
                 "Mohamed Alaa Shebl Mohamed  |  230504583", font_size=16, color=LIGHT_GRAY)
    add_text_box(slide, Inches(1.5), Inches(4.6), Inches(8), Inches(0.5),
                 "Artificial Neural Networks — Istanbul Atlas University", font_size=14, color=DARK_GRAY)
    add_text_box(slide, Inches(1.5), Inches(5.1), Inches(8), Inches(0.5),
                 "April 2025", font_size=14, color=DARK_GRAY)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Problem Statement & Objective", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    # Problem box
    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(2.4), CARD_BG)
    add_text_box(slide, Inches(1.1), Inches(1.5), Inches(5.0), Inches(0.4),
                 "THE PROBLEM", font_size=12, color=ACCENT_TEAL, bold=True)
    add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5.0), Inches(1.6),
                 "Automatically classify handwritten English letters A–Z from "
                 "grayscale images using an Artificial Neural Network built entirely "
                 "from scratch with NumPy — no TensorFlow or PyTorch.",
                 font_size=15, color=LIGHT_GRAY)

    # Objective box
    add_shape_rect(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(3.0), CARD_BG)
    add_text_box(slide, Inches(1.1), Inches(4.2), Inches(5.0), Inches(0.4),
                 "OBJECTIVES", font_size=12, color=ACCENT_TEAL, bold=True)
    add_bullet_list(slide, Inches(1.1), Inches(4.6), Inches(5.0), Inches(2.5), [
        "Build a complete MLP with backpropagation from scratch",
        "Train and evaluate on 762K+ handwritten character images",
        "Provide an interactive GUI for training and prediction",
        "Demonstrate ANN concepts with logic gate examples",
    ], font_size=14, line_spacing=Pt(20))

    # Right side: scope
    add_shape_rect(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.7), CARD_BG)
    add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.2), Inches(0.4),
                 "PROJECT SCOPE", font_size=12, color=ACCENT_TEAL, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.2), Inches(5.0), [
        "Multi-class classification (26 classes: A–Z)",
        "From-scratch NumPy implementation (no ML libraries)",
        "MLP architecture: 784 → 256 → 128 → 26",
        "ReLU hidden activation + Softmax output",
        "Cross-entropy loss + mini-batch gradient descent",
        "Tkinter GUI with live training visualization",
        "Logic gate demonstrations (AND, OR, XOR)",
    ], font_size=14, line_spacing=Pt(18))


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Project Pipeline Overview", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    # Pipeline stages
    stages = [
        ("01", "Dataset\n(.npz)", "762K grayscale\nimages, 28×28"),
        ("02", "Preprocessing", "Normalize, flatten,\nlabel remapping"),
        ("03", "MLP Training", "Forward + Backprop\nmini-batch SGD"),
        ("04", "Evaluation", "Accuracy, loss,\nconfusion matrix"),
        ("05", "GUI\nPrediction", "Draw or upload\nimage to classify"),
    ]

    for i, (num, title, desc) in enumerate(stages):
        x = Inches(0.6 + i * 2.5)
        y = Inches(1.6)
        # Card
        add_shape_rect(slide, x, y, Inches(2.2), Inches(2.8), CARD_BG)
        # Number badge
        add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.6), Inches(0.4),
                     num, font_size=20, color=ACCENT_BLUE, bold=True)
        # Title
        add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(1.9), Inches(0.8),
                     title, font_size=16, color=WHITE, bold=True)
        # Description
        add_text_box(slide, x + Inches(0.15), y + Inches(1.5), Inches(1.9), Inches(1.0),
                     desc, font_size=12, color=DARK_GRAY)
        # Arrow between stages
        if i < len(stages) - 1:
            add_text_box(slide, x + Inches(2.2), y + Inches(1.0), Inches(0.3), Inches(0.4),
                         "→", font_size=20, color=ACCENT_TEAL, bold=True)

    # Project files section
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.4),
                 "KEY PROJECT FILES", font_size=12, color=ACCENT_TEAL, bold=True)

    files = [
        ("src/ann/mlp.py", "Core MLP neural network implementation"),
        ("src/ann/data.py", "Dataset loading, preprocessing, train/test split"),
        ("src/ann/utils.py", "One-hot encoding, metrics, confusion matrix"),
        ("train_mlp.py", "Command-line training and evaluation"),
        ("gui_app.py", "Interactive Tkinter GUI application"),
    ]

    for i, (fname, fdesc) in enumerate(files):
        y = Inches(5.2 + i * 0.4)
        add_text_box(slide, Inches(1.0), y, Inches(2.5), Inches(0.4),
                     fname, font_size=13, color=ACCENT_BLUE, bold=True,
                     font_name="Consolas")
        add_text_box(slide, Inches(3.8), y, Inches(5), Inches(0.4),
                     fdesc, font_size=13, color=LIGHT_GRAY)

    # Add architecture diagram if available
    add_image_safe(slide, OUTPUT_DIR / "neural_network_training_flow.png",
                   Inches(7.5), Inches(4.8), width=Inches(5.2))


def slide_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Dataset Overview", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    # Stats cards
    add_stat_card(slide, Inches(0.8), Inches(1.5), "762,213", "Total Images", ACCENT_BLUE)
    add_stat_card(slide, Inches(3.7), Inches(1.5), "28 × 28", "Image Size (px)", ACCENT_TEAL)
    add_stat_card(slide, Inches(6.6), Inches(1.5), "26", "Classes (A–Z)", ACCENT_ORANGE)
    add_stat_card(slide, Inches(9.5), Inches(1.5), "Grayscale", "Color Mode", ACCENT_BLUE)

    # Dataset details
    add_shape_rect(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(3.5), CARD_BG)
    add_text_box(slide, Inches(1.1), Inches(3.6), Inches(5), Inches(0.4),
                 "DATASET DETAILS", font_size=12, color=ACCENT_TEAL, bold=True)
    add_bullet_list(slide, Inches(1.1), Inches(4.0), Inches(5.0), Inches(2.8), [
        "Source: Kaggle Alphabet Characters Fonts Dataset",
        "Includes handwritten character samples",
        "Stored as NumPy .npz format (images + labels)",
        "Class imbalance: min 16,109 — max 72,816 per class",
        "Stratified split preserves class proportions",
        "Split: 80% train / 20% test (seed=42)",
    ], font_size=14, line_spacing=Pt(18))

    # Preprocessing
    add_shape_rect(slide, Inches(6.7), Inches(3.5), Inches(5.8), Inches(3.5), CARD_BG)
    add_text_box(slide, Inches(7.0), Inches(3.6), Inches(5.2), Inches(0.4),
                 "PREPROCESSING PIPELINE", font_size=12, color=ACCENT_TEAL, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(4.0), Inches(5.2), Inches(2.8), [
        "Resize to 28×28 (bilinear interpolation)",
        "Normalize pixels: [0, 255] → [0.0, 1.0]",
        "Flatten: 28×28 → 784-element vector",
        "Remap labels to contiguous 0–25",
        "GUI: auto-crop, center, autocontrast",
        "GUI: dual polarity (normal + inverted)",
    ], font_size=14, line_spacing=Pt(18))


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "MLP Architecture & Design", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    # Architecture visualization
    layers = [
        ("Input", "784", "28×28 pixels\nflattened"),
        ("Hidden 1", "256", "ReLU\nactivation"),
        ("Hidden 2", "128", "ReLU\nactivation"),
        ("Output", "26", "Softmax\n(A–Z classes)"),
    ]

    for i, (name, size, desc) in enumerate(layers):
        x = Inches(0.8 + i * 3.1)
        # Layer box
        accent = ACCENT_BLUE if i < 3 else ACCENT_ORANGE
        add_shape_rect(slide, x, Inches(1.5), Inches(2.6), Inches(2.2), CARD_BG, accent, Pt(2))
        add_text_box(slide, x + Inches(0.2), Inches(1.6), Inches(2.2), Inches(0.4),
                     name, font_size=14, color=accent, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(2.2), Inches(0.6),
                     f"{size} neurons", font_size=24, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.2), Inches(0.6),
                     desc, font_size=12, color=DARK_GRAY)

        if i < len(layers) - 1:
            add_text_box(slide, x + Inches(2.6), Inches(2.2), Inches(0.5), Inches(0.4),
                         "→", font_size=24, color=ACCENT_TEAL, bold=True)

    # Key formulas
    add_shape_rect(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(3.0), CARD_BG)
    add_text_box(slide, Inches(1.1), Inches(4.2), Inches(5), Inches(0.4),
                 "KEY EQUATIONS", font_size=12, color=ACCENT_TEAL, bold=True)
    add_bullet_list(slide, Inches(1.1), Inches(4.6), Inches(5.0), Inches(2.3), [
        "Forward:  z = a·W + b,  a = ReLU(z)",
        "Softmax:  P(class_i) = e^(z_i) / Σ e^(z_j)",
        "Loss:  L = -(1/m) Σ y·log(ŷ)  (cross-entropy)",
        "Backprop:  δ = ŷ - y  (output layer)",
        "Update:   W := W - η · ∇W",
    ], font_size=14, line_spacing=Pt(16))

    # Parameters card
    add_shape_rect(slide, Inches(6.7), Inches(4.1), Inches(5.8), Inches(3.0), CARD_BG)
    add_text_box(slide, Inches(7.0), Inches(4.2), Inches(5.2), Inches(0.4),
                 "PARAMETER COUNT", font_size=12, color=ACCENT_TEAL, bold=True)

    params = [
        ("Layer 1 (784→256)", "200,960"),
        ("Layer 2 (256→128)", "32,896"),
        ("Layer 3 (128→26)", "3,354"),
        ("Total Parameters", "237,210"),
    ]
    for i, (label, val) in enumerate(params):
        y = Inches(4.7 + i * 0.55)
        color = WHITE if i == len(params) - 1 else LIGHT_GRAY
        bold = i == len(params) - 1
        add_text_box(slide, Inches(7.2), y, Inches(3.0), Inches(0.4),
                     label, font_size=14, color=color, bold=bold)
        add_text_box(slide, Inches(10.5), y, Inches(1.5), Inches(0.4),
                     val, font_size=14, color=ACCENT_BLUE, bold=bold,
                     alignment=PP_ALIGN.RIGHT)


def slide_training(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Training Configuration", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    # Configuration table
    configs = [
        ("Learning Rate", "0.01"),
        ("Batch Size", "256"),
        ("Max Epochs", "20"),
        ("Early Stopping", "Patience = 5"),
        ("Validation Split", "10% of training data"),
        ("Test Split", "20% of full dataset"),
        ("Activation (Hidden)", "ReLU"),
        ("Activation (Output)", "Softmax"),
        ("Initialization", "He (ReLU) / Xavier (Sigmoid)"),
        ("Optimizer", "Mini-batch Gradient Descent"),
        ("Random Seed", "42"),
    ]

    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.7), CARD_BG)
    add_text_box(slide, Inches(1.1), Inches(1.5), Inches(5), Inches(0.4),
                 "HYPERPARAMETERS", font_size=12, color=ACCENT_TEAL, bold=True)

    for i, (param, val) in enumerate(configs):
        y = Inches(1.95 + i * 0.45)
        bg_color = SECTION_BG if i % 2 == 0 else CARD_BG
        add_shape_rect(slide, Inches(1.0), y, Inches(5.0), Inches(0.4), bg_color)
        add_text_box(slide, Inches(1.1), y, Inches(2.6), Inches(0.4),
                     param, font_size=13, color=LIGHT_GRAY)
        add_text_box(slide, Inches(3.7), y, Inches(2.2), Inches(0.4),
                     val, font_size=13, color=WHITE, bold=True)

    # Training features
    add_shape_rect(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.7), CARD_BG)
    add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.2), Inches(0.4),
                 "TRAINING FEATURES", font_size=12, color=ACCENT_TEAL, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.8), [
        "Mini-batch processing for memory efficiency",
        "Data shuffling each epoch to prevent ordering bias",
        "Epoch-level loss/accuracy tracking",
        "Validation monitoring after each epoch",
        "Early stopping with best-weights restoration",
        "Reproducible splits with fixed random seeds",
        "Epoch callbacks for GUI live updates",
        "Background thread training (GUI stays responsive)",
    ], font_size=14, line_spacing=Pt(22))


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Results & Evaluation", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    # Big stats
    add_stat_card(slide, Inches(0.8), Inches(1.4), "78.15%", "Test Accuracy", ACCENT_TEAL)
    add_stat_card(slide, Inches(3.7), Inches(1.4), "0.8996", "Test Loss", ACCENT_ORANGE)
    add_stat_card(slide, Inches(6.6), Inches(1.4), "152,452", "Test Samples", ACCENT_BLUE)
    add_stat_card(slide, Inches(9.5), Inches(1.4), "237,210", "Parameters", ACCENT_TEAL)

    # Training curve image
    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(5), Inches(0.4),
                 "TRAINING CURVES", font_size=12, color=ACCENT_TEAL, bold=True)
    add_image_safe(slide, OUTPUT_DIR / "mlp_training_curve.png",
                   Inches(0.8), Inches(3.7), width=Inches(6.0))


def slide_confusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Confusion Matrix & Error Analysis", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    # Confusion matrix image
    add_image_safe(slide, OUTPUT_DIR / "mlp_confusion_matrix.png",
                   Inches(0.5), Inches(1.4), height=Inches(5.5))

    # Error analysis
    add_shape_rect(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.7), CARD_BG)
    add_text_box(slide, Inches(7.3), Inches(1.5), Inches(5.2), Inches(0.4),
                 "TOP CONFUSIONS", font_size=12, color=ACCENT_TEAL, bold=True)

    confusions = [
        ("D → O", "474", "Similar round shapes"),
        ("Y → V", "336", "Similar angular form"),
        ("F → P", "317", "Similar top structure"),
        ("I → J", "315", "Similar vertical strokes"),
        ("Q → O", "312", "Q without tail = O"),
        ("H → N", "298", "Similar stroke patterns"),
        ("E → C", "280", "Similar curved shape"),
    ]

    for i, (pair, count, reason) in enumerate(confusions):
        y = Inches(2.0 + i * 0.65)
        add_text_box(slide, Inches(7.3), y, Inches(1.2), Inches(0.4),
                     pair, font_size=14, color=ACCENT_ORANGE, bold=True,
                     font_name="Consolas")
        add_text_box(slide, Inches(8.6), y, Inches(0.8), Inches(0.4),
                     count, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, Inches(9.5), y, Inches(3.0), Inches(0.4),
                     reason, font_size=12, color=LIGHT_GRAY)

    add_text_box(slide, Inches(7.3), Inches(6.5), Inches(5.0), Inches(0.5),
                 "Visually similar letters are the primary source of errors.",
                 font_size=13, color=DARK_GRAY)


def slide_logic_gates(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Bonus: Logic Gate Demonstrations", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    gates = [
        ("AND Gate", "logic_gate_and.png", "Linearly Separable ✓"),
        ("OR Gate", "logic_gate_or.png", "Linearly Separable ✓"),
        ("XOR Gate", "logic_gate_xor.png", "NOT Linearly Separable ✗"),
    ]

    for i, (title, img, note) in enumerate(gates):
        x = Inches(0.5 + i * 4.2)
        add_text_box(slide, x, Inches(1.3), Inches(3.8), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_image_safe(slide, OUTPUT_DIR / img,
                       x, Inches(1.8), width=Inches(3.8))
        note_color = ACCENT_TEAL if "✓" in note else ACCENT_ORANGE
        add_text_box(slide, x, Inches(5.6), Inches(3.8), Inches(0.4),
                     note, font_size=13, color=note_color, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.8),
                 "XOR demonstrates why hidden layers are essential — "
                 "a single-layer perceptron cannot solve non-linearly separable problems.",
                 font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_gui(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Interactive GUI Application", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    # GUI features - left
    add_shape_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5), CARD_BG)
    add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.4),
                 "GUI FEATURES", font_size=12, color=ACCENT_TEAL, bold=True)
    add_bullet_list(slide, Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.8), [
        "Browse and select any .npz dataset file",
        "Configure all hyperparameters from the interface",
        "Train model with live loss/accuracy curve updates",
        "Drawing canvas (320×320) with adjustable brush",
        "Upload external images for classification",
        "Predicted class with confidence percentage",
        "Top-k predictions displayed in table format",
        "Probability bar chart visualization",
        "Background thread training (GUI stays responsive)",
        "Model save/load functionality",
    ], font_size=14, line_spacing=Pt(16))

    # Technical details - right
    add_shape_rect(slide, Inches(6.7), Inches(1.5), Inches(5.8), Inches(5.5), CARD_BG)
    add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.2), Inches(0.4),
                 "TECHNICAL DESIGN", font_size=12, color=ACCENT_TEAL, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.8), [
        "Built with Python Tkinter + Matplotlib embedding",
        "Thread-safe queue-based event system",
        "Worker thread for training, main thread for UI",
        "Event polling via Tkinter after() every 80ms",
        "PIL-based image preprocessing pipeline",
        "Auto-detects dark/light background polarity",
        "Dual prediction: tries normal + inverted input",
        "Graceful error handling with user-friendly messages",
        "Status log with automatic line trimming",
        "Proper cleanup on window close",
    ], font_size=14, line_spacing=Pt(16))


def slide_gui_screenshots(prs):
    """2×2 grid of GUI mockup screenshots."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
                 "GUI Screenshots — Application in Action", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    screenshots = [
        ("gui_main.png",             "Main Window — Draw & Predict Tab",    ACCENT_BLUE),
        ("gui_prediction.png",       "After Prediction — Letter 'A' Drawn",  ACCENT_TEAL),
        ("gui_training_metrics.png", "Training Metrics — Live Loss/Accuracy",ACCENT_ORANGE),
        ("gui_probability.png",      "Probability View — Top Class Chart",   ACCENT_TEAL),
    ]

    positions = [
        (Inches(0.3),  Inches(1.3)),   # top-left
        (Inches(6.85), Inches(1.3)),   # top-right
        (Inches(0.3),  Inches(4.35)),  # bottom-left
        (Inches(6.85), Inches(4.35)),  # bottom-right
    ]

    img_w = Inches(6.3)
    img_h = Inches(2.85)

    for (img_file, caption, accent), (lx, ly) in zip(screenshots, positions):
        img_path = OUTPUT_DIR / img_file
        # Thin accent border frame
        add_shape_rect(slide, lx - Inches(0.05), ly - Inches(0.05),
                       img_w + Inches(0.10), img_h + Inches(0.10),
                       CARD_BG, accent, Pt(1.5))
        add_image_safe(slide, img_path, lx, ly, width=img_w)
        # Caption below image
        add_text_box(slide, lx, ly + img_h + Inches(0.04),
                     img_w, Inches(0.28),
                     caption, font_size=11, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER)


def slide_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Limitations & Future Work", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    # Limitations
    add_shape_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5), CARD_BG)
    add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.4),
                 "CURRENT LIMITATIONS", font_size=12, color=ACCENT_ORANGE, bold=True)
    add_bullet_list(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.5), [
        "MLP on raw pixels — no spatial feature extraction",
        "CNNs significantly outperform MLPs on image tasks",
        "Class imbalance biases predictions toward common classes",
        "No data augmentation (rotation, shift, noise)",
        "Fixed learning rate throughout training",
        "No regularization (dropout, L2 penalty)",
    ], font_size=14, bullet_color=ACCENT_ORANGE, line_spacing=Pt(24))

    # Future work
    add_shape_rect(slide, Inches(6.7), Inches(1.5), Inches(5.8), Inches(5.5), CARD_BG)
    add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.2), Inches(0.4),
                 "FUTURE IMPROVEMENTS", font_size=12, color=ACCENT_TEAL, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(2.1), Inches(5.2), Inches(4.5), [
        "Upgrade to CNN architecture for spatial features",
        "Add data augmentation pipeline (rotations, shifts)",
        "Apply class weighting or focal loss for imbalance",
        "Implement learning rate scheduling / Adam optimizer",
        "Add dropout and L2 regularization",
        "Extend to lowercase letters, digits, and symbols",
    ], font_size=14, line_spacing=Pt(24))


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Conclusion", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(3.0), ACCENT_BLUE)

    achievements = [
        ("From-Scratch ANN", "Complete MLP implementation using only NumPy —\nforward pass, backpropagation, softmax, cross-entropy, mini-batch SGD"),
        ("78.15% Accuracy", "Trained on 762K images and evaluated on 152K test samples\nwith reproducible stratified splits"),
        ("Interactive GUI", "Full Tkinter application with live training plots,\ndrawing canvas, image upload, and probability visualization"),
        ("ANN Concepts Demo", "Logic gate examples (AND, OR, XOR) illustrate\nlinear separability and the need for hidden layers"),
    ]

    for i, (title, desc) in enumerate(achievements):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(1.5 + (i // 2) * 2.7)
        add_shape_rect(slide, x, y, Inches(5.7), Inches(2.2), CARD_BG)
        accent = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_ORANGE, ACCENT_BLUE][i]
        add_accent_line(slide, x + Inches(0.3), y + Inches(0.2), Inches(2.5), accent)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.4), Inches(5.0), Inches(0.5),
                     title, font_size=20, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(5.0), Inches(1.0),
                     desc, font_size=14, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
                 "This project demonstrates a strong understanding of ANN fundamentals "
                 "through a complete, functional, and well-documented implementation.",
                 font_size=16, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)


def slide_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_accent_line(slide, Inches(4.5), Inches(2.0), Inches(4.3), ACCENT_TEAL)

    add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.0),
                 "Thank You", font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.6),
                 "Questions & Discussion", font_size=22, color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(4.3), Inches(2.3), CARD_BG)

    add_text_box(slide, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.5),
                 "Mohamed Alaa Shebl Mohamed  |  230504583", font_size=16, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.5),
                 "Artificial Neural Networks — Istanbul Atlas University",
                 font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ── Main ────────────────────────────────────────────────────────────

def main():
    prs = Presentation()

    # Set widescreen 16:9
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)             # 1
    slide_problem(prs)           # 2
    slide_pipeline(prs)          # 3
    slide_dataset(prs)           # 4
    slide_architecture(prs)      # 5
    slide_training(prs)          # 6
    slide_results(prs)           # 7
    slide_confusion(prs)         # 8
    slide_logic_gates(prs)       # 9
    slide_gui(prs)               # 10
    slide_gui_screenshots(prs)   # 11  ← NEW: GUI screenshots
    slide_limitations(prs)       # 12
    slide_conclusion(prs)        # 13
    slide_thank_you(prs)         # 14

    prs.save(str(PPTX_PATH))
    print(f"Presentation saved to: {PPTX_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
